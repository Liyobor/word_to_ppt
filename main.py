from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RGB

def read_word_file(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def format_paragraph_text(paragraph_text):
    # Split the text using spaces and rejoin with newline to separate the two parts
    parts = paragraph_text.split(maxsplit=1)
    return '\n'.join(parts)

def calculate_string_length(s):
    length = 0
    for char in s:
        if '\u0021' <= char <= '\u007E':  # if the character is a Latin character (excluding space)
            length += 1
        else:
            length += 2
    return length


file_path = 'sunnyday.docx'
text = read_word_file(file_path)
prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height
for paragraph in text.split("\n\n"):
    if paragraph:  # 確保段落不是空的
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用 "Blank" 的版面設計

        # 設定背景色為黑色
        # background = slide.background
        # fill = background.fill
        # fill.solid()
        # fill.fore_color.rgb = RGB(0, 0, 0)

        # txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
        # tf = txBox.text_frame
        shape = slide.shapes.add_shape(
                1,  # Rectangle shape type
                0, 0, slide_width, slide_height
            )
                    # Set the shape fill to black
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGB(0, 0, 0)
        tf = shape.text_frame
        tf.clear()  # remove any existing paragraphs
        
        p = tf.add_paragraph()
        # formatted_text = format_paragraph_text(paragraph)
        
        p.text = paragraph
        indices = [i for i, char in enumerate(paragraph) if char == "\n"]
        
        
        
        print(f'paragraph = {paragraph}')

        # print(f'indices = {indices}')
        max_length = 0

        
        for i in range(len(indices)):
            if i == 0:
                font_length = calculate_string_length(paragraph[:indices[i]])
            else:
                font_length = calculate_string_length(paragraph[indices[i-1]:indices[i]])
            max_length = max(max_length,font_length)
        if max_length==0:
            print("無換行")
            max_length = calculate_string_length(paragraph)

        # max_length
        # print(f'slide_width = {slide_width}')
        # print(f"max_length = {max_length}")
        
        # print(f'換行 = {indices}')
        
        # print(f'p text = {paragraph.split("/n")}')
        # 根據段落的文字量動態調整字體大小
        # print(len(paragraph[0:paragraph.find('\n')]))
        # print('---')
        '''
        字體5時,可以容納284個拉丁字母長度
        因此字體為n*5時,可容納284/n的拉丁字母長度
        '''
        
        print(f'max len = {max_length}')
        print(f'font size = {int((5*284/(max_length)))}')
        # font_size = 284/max_length
        font_size = Pt(int((5*283/(max_length))) - 3) # 避免誤差
        # if max_length > 40:
        #     font_size = Pt(30)
        # elif max_length > 18:
        #     font_size = Pt(40)
        # else:
        #     font_size = Pt(60)
        

        # print(f'font size = {font_size}')
        # 設定文字的樣式
        for run in p.runs:
            font = run.font
            font.bold = True
            font.size = font_size
            font.color.rgb = RGB(255, 255, 255)
        p.alignment = 1  # center alignment

        

prs.save("example.pptx")
